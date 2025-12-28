from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

MAX_CHARS_PER_SLIDE = 800

def split_text_into_chunks(text, max_chars=MAX_CHARS_PER_SLIDE):
    words = text.split()
    chunks, chunk = [], ""

    for w in words:
        if len(chunk) + len(w) + 1 <= max_chars:
            chunk += (" " if chunk else "") + w
        else:
            chunks.append(chunk)
            chunk = w
    if chunk:
        chunks.append(chunk)

    return chunks

def format_summary_to_bullets(summary, max_bullets=6):
    sentences = [s.strip() for s in summary.replace("\n", " ").split(".") if len(s.strip()) > 30]
    bullets = sentences[:max_bullets]
    return bullets

def generate_slides(papers, filename="AI_Research_Summary.pptx"):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[1]  # Title + Content

    for paper in papers:
        title = paper["title"]

        base_info = (
            f"Authors: {', '.join(paper['authors'])}\n"
            f"Published: {paper['published']}\n"
            f"PDF: {paper['pdf_url']}\n\n"
        )

        full_text = base_info + paper["summary"]
        chunks = split_text_into_chunks(full_text)

        for chunk in chunks:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title

            tf = slide.placeholders[1].text_frame
            tf.clear()

            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            bullets = format_summary_to_bullets(chunk)

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 1

                p.font.size = Pt(20)
                p.font.bold = False

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer